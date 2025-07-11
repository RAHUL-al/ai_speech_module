import os
import ujson as json
import redis
from datetime import datetime, timedelta
from fastapi import FastAPI, Depends, HTTPException, status, Request, WebSocket, WebSocketDisconnect,UploadFile, File, Form
from fastapi_sqlalchemy import DBSessionMiddleware, db
from fastapi.middleware.cors import CORSMiddleware
from pydub import AudioSegment
from models import User, Essay
from schemas import UserCreate, UserOut, UserUpdate, Token, LoginRequest, ForgotPasswordRequest, GeminiRequest, TextToSpeechRequest
from ai_speech_module import Topic as AiTopic
from auth import hash_password, verify_password, create_access_token
from dotenv import load_dotenv
from urllib.parse import parse_qs
from concurrent.futures import ThreadPoolExecutor
import asyncio
from fastapi.responses import FileResponse
import time
from google.cloud import vision
import pinecone
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import Pinecone
from langchain.llms import OpenAI
from langchain.chains import RetrievalQA
from typing import List

executor = ThreadPoolExecutor(max_workers=4)

load_dotenv()

app = FastAPI(title="FastAPI​‑SQLAlchemy​‑MySQL")

redis_client = redis.StrictRedis(
    host=os.getenv("REDIS_HOST"),
    port=int(os.getenv("REDIS_PORT")),
    username=os.getenv("REDIS_USERNAME"),
    password=os.getenv("REDIS_PASSWORD"),
    decode_responses=True
)

origins = [
    "http://localhost:5173",
    "https://02d59704d602.ngrok-free.app",
]
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

app.add_middleware(DBSessionMiddleware, db_url=os.environ["DATABASE_URL"])

def get_user_from_redis_session(request: Request):
    token = request.headers.get("Authorization")
    if not token or not token.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid token")
    token = token.split(" ")[1]
    session_data = redis_client.get(f"session:{token}")
    if not session_data:
        raise HTTPException(status_code=401, detail="Session expired or invalid")
    return json.loads(session_data)

@app.post("/register", response_model=UserOut, status_code=status.HTTP_201_CREATED)
def register(user: UserCreate):
    if db.session.query(User).filter(
        (User.username == user.username) | (User.email == user.email)
    ).first():
        raise HTTPException(400, "Username or email already exists")

    new_user = User(
        username=user.username,
        email=user.email,
        password=hash_password(user.password)
    )
    db.session.add(new_user)
    db.session.commit()
    db.session.refresh(new_user)
    return new_user

@app.post("/login", response_model=Token)
def login(data: LoginRequest):
    user = db.session.query(User).filter(User.username == data.username).first()
    if not user or not verify_password(data.password, user.password):
        raise HTTPException(status.HTTP_401_UNAUTHORIZED, detail="Bad credentials")

    token = create_access_token({"sub": str(user.id)})
    session_payload = json.dumps({"user_id": user.id, "username": user.username})
    redis_client.setex(f"session:{token}", timedelta(hours=1), session_payload)

    return Token(access_token=token, username=user.username)

@app.get("/logout")
def logout(request: Request):
    token = request.headers.get("Authorization")
    if token and token.startswith("Bearer "):
        redis_client.delete(f"session:{token.split(' ')[1]}")
    return {"detail": "Logged out"}

@app.get("/users/{user_id}", response_model=UserOut)
def get_user(user_id: int, user=Depends(get_user_from_redis_session)):
    user = db.session.get(User, user_id)
    if not user:
        raise HTTPException(404, "User not found")
    return user

@app.get("/me", response_model=UserOut)
def me(user=Depends(get_user_from_redis_session)):
    return user

@app.put("/users/{user_id}", response_model=UserOut)
def update_user(user_id: int, payload: UserUpdate, user=Depends(get_user_from_redis_session)):
    user_obj = db.session.get(User, user_id)
    if not user_obj:
        raise HTTPException(404, "User not found")
    if payload.username: user_obj.username = payload.username
    if payload.email: user_obj.email = payload.email
    if payload.password: user_obj.password = hash_password(payload.password)
    db.session.commit()
    return user_obj

@app.delete("/users/{user_id}", status_code=status.HTTP_204_NO_CONTENT)
def delete_user(user_id: int, user=Depends(get_user_from_redis_session)):
    user = db.session.get(User, user_id)
    if not user:
        raise HTTPException(404, "User not found")
    db.session.delete(user)
    db.session.commit()
    return None

@app.post("/forgot-password", status_code=status.HTTP_200_OK)
def forgot_password(request: ForgotPasswordRequest):
    user = db.session.query(User).filter(User.username == request.username).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    user.password = hash_password(request.new_password)
    db.session.commit()
    return {"detail": "Password reset successfully"}

@app.post("/generate-prompt")
def generate_prompt(data: GeminiRequest, user=Depends(get_user_from_redis_session)):
    prompt = (
        f"Generate a essay for a student in class {data.student_class} "
        f"with a {data.accent} accent, on the topic '{data.topic}', and the mood is '{data.mood}' and give me essay less than 400 words."
    )
    username = user.get("username")
    topic = AiTopic()
    response_text = topic.topic_data_model_for_Qwen(prompt, username=username)

    essay = Essay(
        student_class=data.student_class,
        accent=data.accent,
        topic=data.topic,
        mood=data.mood,
        content=response_text,
        user_id=user["user_id"]
    )
    db.session.add(essay)
    db.session.commit()
    db.session.refresh(essay)

    return {"response": response_text, "essay_id": essay.id}

TEMP_DIR = os.path.abspath("audio_folder")
os.makedirs(TEMP_DIR, exist_ok=True)

@app.websocket("/ws/audio")
async def audio_ws(websocket: WebSocket):
    await websocket.accept()
    query_params = parse_qs(websocket.url.query)
    username = query_params.get("username", [None])[0]

    if not username:
        await websocket.close(code=4001)
        print("Username not provided in WebSocket connection.")
        return

    print(f"[WS] Client connected: {username}")
    chunk_index = 0
    chunk_files = []
    text_output = []
    topic = AiTopic()

    date_str = datetime.now().strftime("%Y-%m-%d")
    user_dir = os.path.join(TEMP_DIR, username, date_str)
    os.makedirs(user_dir, exist_ok=True)

    final_output = os.path.join(user_dir, f"{username}_output.wav")
    transcript_path = os.path.join(user_dir, f"{username}_transcript.txt")

    if os.path.exists(final_output):
        os.remove(final_output)
    if os.path.exists(transcript_path):
        os.remove(transcript_path)

    loop = asyncio.get_event_loop()

    try:
        while True:
            message = await websocket.receive()

            if message["type"] == "websocket.disconnect":
                print(f"[WS] {username} disconnected.")
                break

            if message["type"] == "websocket.receive" and "bytes" in message:
                chunk_filename = os.path.join(user_dir, f"chunk_{chunk_index}.wav")
                audio = AudioSegment(
                    data=message["bytes"],
                    sample_width=2,
                    frame_rate=16000,
                    channels=1
                )
                audio.export(chunk_filename, format="wav")
                chunk_files.append(chunk_filename)

                # Real-time parallel processing
                results = await asyncio.gather(
                    asyncio.to_thread(topic.speech_to_text, chunk_filename),
                    asyncio.to_thread(topic.analyze_emotion, chunk_filename),
                    asyncio.to_thread(topic.analyze_fluency, chunk_filename),
                    asyncio.to_thread(topic.analyze_pronunciation, chunk_filename),
                )

                transcribed_text, emotion, fluency, pronunciation = results

                text_output.append(transcribed_text)
                print(f"[Chunk {chunk_index}] Transcribed: {transcribed_text.strip()}")
                print(f"Emotion: {emotion} | Fluency: {fluency} | Pronunciation: {pronunciation}")
                chunk_index += 1

    except WebSocketDisconnect:
        print(f"[WS] {username} forcibly disconnected.")

    finally:
        await loop.run_in_executor(None, merge_chunks, chunk_files, final_output)
        print(f"[Output] Final audio saved: {final_output}")

        with open(transcript_path, "w", encoding="utf-8") as f:
            f.write(" ".join(text_output).strip())
        print(f"[Output] Transcript saved: {transcript_path}")

        for file in chunk_files:
            try:
                os.remove(file)
            except Exception as e:
                print(f"[Warning] Failed to remove {file}: {e}")

async def merge_chunks(chunk_files, final_output):
    print("[Merge] Merging audio chunks...")
    combined = AudioSegment.empty()
    for file in chunk_files:
        audio = AudioSegment.from_file(file, format="wav")
        combined += audio
    combined.export(final_output, format="wav")
    print("[Merge] Merged audio file saved.")

@app.get("/grammar-score")
def grammar_score(username: str):
    try:
        service = AiTopic()
        result = service.overall_scoring(username=username)
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/get-tts-audio")
def get_tts_audio(username: str):
    folder = os.path.join("text_to_speech_audio_folder", username)
    file_path = os.path.join(folder, f"{username}_output.wav")

    timeout = 60
    poll_interval = 2
    waited = 0

    while waited < timeout:
        if os.path.exists(file_path):
            return FileResponse(
                file_path,
                media_type="audio/wav",
                filename=f"{username}_output.wav"
            )
        time.sleep(poll_interval)
        waited += poll_interval

    raise HTTPException(status_code=408, detail="Audio file not generated within 1 minute.")


@app.get("/test")
def welcome_page():
    return {"Message": "Welcome the ai speech module page."}

@app.get("/")
def root():
    return {"message": "API is working"}


pinecone.init(api_key=os.getenv("PINECONE_API_KEY"), environment=os.getenv("PINECONE_ENV"))
INDEX_NAME = "teacher_rag"
if INDEX_NAME not in pinecone.list_indexes():
    pinecone.create_index(INDEX_NAME, dimension=1536)

index = pinecone.Index(INDEX_NAME)
embedding_fn = OpenAIEmbeddings()

vision_client = vision.ImageAnnotatorClient()

def ocr_image(file_bytes: bytes) -> str:
    image = vision.Image(content=file_bytes)
    response = vision_client.document_text_detection(image=image)
    return response.full_text_annotation.text or ""

def extract_text_from_pdf(file_bytes: bytes) -> str:
    from pdfminer.high_level import extract_text_to_fp
    import io
    out = io.StringIO()
    extract_text_to_fp(io.BytesIO(file_bytes), out)
    return out.getvalue()

def extract_text_from_ppt(file_bytes: bytes) -> str:
    from pptx import Presentation
    import io
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

@app.post("/upload-documents")
async def upload_documents(
    class_name: str = Form(...),
    subject: str = Form(...),
    curriculum: str = Form(...),
    teacher_id: str = Form(...),
    files: List[UploadFile] = File(...)
):
    all_chunks = []
    metadata_list = []
    for file in files:
        content = await file.read()
        ext = os.path.splitext(file.filename)[1].lower()
        if ext in [".pdf"]:
            text = extract_text_from_pdf(content)
        elif ext in [".pptx", ".ppt"]:
            text = extract_text_from_ppt(content)
        elif ext in [".jpg", ".jpeg", ".png"]:
            text = ocr_image(content)
        else:
            continue
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_text(text)
        for chunk in chunks:
            metadata_list.append({
                "teacher_id": teacher_id,
                "class": class_name,
                "subject": subject,
                "curriculum": curriculum,
                "source": file.filename
            })
            all_chunks.append(chunk)
    if not all_chunks:
        raise HTTPException(400, "No valid content extracted")
    embeddings = embedding_fn.embed_documents(all_chunks)
    Pinecone.from_texts(all_chunks, embedding_fn,
                       index_name=INDEX_NAME, metadatas=metadata_list)
    return {"message": f"Indexed {len(all_chunks)} chunks for teacher {teacher_id}"}



@app.post("/rag-chat")
async def rag_chat(
    question: str,
    teacher_id: str = Form(...),
    class_name: str = Form(None),
    subject: str = Form(None)
):
    filter_meta = {"teacher_id": teacher_id}
    if class_name: filter_meta["class"] = class_name
    if subject: filter_meta["subject"] = subject

    vectordb = Pinecone.from_existing_index(index_name=INDEX_NAME,
                                             embedding=embedding_fn,
                                             filter=filter_meta)
    llm = OpenAI(temperature=0, model="gpt-3.5-turbo")
    qa = RetrievalQA.from_chain_type(llm=llm, chain_type="stuff", retriever=vectordb.as_retriever())
    answer = qa.run(question)
    return {"answer": answer}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8888, reload=True)